"""
Build a 10-slide research presentation as PowerPoint (.pptx).
Slide deck: concise, visual, structured for academic presentation.
"""

import io, json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

# ── Paths ──────────────────────────────────────────────────────────────────
BASE    = Path("C:/Users/aggarwm1/Videos/multi-disciplinary_Hyperspectral")
OUT_DIR = BASE / "outputs"
FIG_DIR = OUT_DIR / "figures"
PPTX_OUT = BASE / "model" / "Presentation_Hyperspectral_Forest.pptx"

# ── Load results ───────────────────────────────────────────────────────────
with open(OUT_DIR / "baseline_multisite_casi_loso/experiment_summary.json") as f:
    baseline = json.load(f)
with open(OUT_DIR / "ablations/als_fusion.json") as f:
    als = json.load(f)
with open(OUT_DIR / "ablations/cross_sensor_bk.json") as f:
    cross = json.load(f)

# ── Colour palette ─────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1F, 0x49, 0x7D)
MED_BLUE    = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE  = RGBColor(0xBD, 0xD7, 0xEE)
ACCENT      = RGBColor(0xED, 0x7D, 0x31)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY   = RGBColor(0x40, 0x40, 0x40)
LIGHT_GREY  = RGBColor(0xF2, 0xF2, 0xF2)

# ── Slide dimensions: 16:9 widescreen ─────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


# ══════════════════════════════════════════════════════════════════════════
# Helpers
# ══════════════════════════════════════════════════════════════════════════

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def rgb_hex(r, g, b):
    return RGBColor(r, g, b)

def fill_slide_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_para_in_tf(tf, text, size=14, bold=False, color=DARK_GREY,
                   align=PP_ALIGN.LEFT, space_before=6, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def add_image_stream(slide, stream, l, t, w, h=None):
    if h:
        slide.shapes.add_picture(stream, l, t, w, h)
    else:
        slide.shapes.add_picture(stream, l, t, w)

def add_image_path(slide, path, l, t, w, h=None):
    if h:
        slide.shapes.add_picture(str(path), l, t, w, h)
    else:
        slide.shapes.add_picture(str(path), l, t, w)

def header_bar(slide, title, subtitle=None):
    """Dark blue header band at top of slide."""
    add_rect(slide, 0, 0, W, Inches(1.15), DARK_BLUE)
    add_text(slide, title,
             Inches(0.35), Inches(0.08), Inches(10), Inches(0.65),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.35), Inches(0.72), Inches(11), Inches(0.38),
                 size=14, bold=False, color=LIGHT_BLUE)
    # slide number placeholder – orange accent bar on right
    add_rect(slide, W - Inches(1.2), 0, Inches(1.2), Inches(1.15), ACCENT)

def slide_number(slide, n):
    add_text(slide, str(n),
             W - Inches(1.1), Inches(0.3), Inches(0.9), Inches(0.5),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def fig_to_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=180, bbox_inches="tight",
                facecolor="none", transparent=True)
    buf.seek(0)
    plt.close(fig)
    return buf

def bullet_box(slide, items, l, t, w, h,
               bg_color=LIGHT_GREY, size=13.5, title=None, title_color=DARK_BLUE):
    add_rect(slide, l, t, w, h, bg_color)
    txBox = slide.shapes.add_textbox(l + Inches(0.18), t + Inches(0.14),
                                      w - Inches(0.3), h - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = title
        r.font.size = Pt(size + 1); r.font.bold = True
        r.font.color.rgb = title_color
        first = False
    for item in items:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        r = p.add_run(); r.text = "•  " + item
        r.font.size = Pt(size)
        r.font.color.rgb = DARK_GREY
        first = False

def stat_box(slide, value, label, l, t, w=Inches(2.4), h=Inches(1.3),
             bg=MED_BLUE, val_color=WHITE, lbl_color=LIGHT_BLUE, val_size=32):
    add_rect(slide, l, t, w, h, bg)
    add_text(slide, value, l, t + Inches(0.1), w, Inches(0.75),
             size=val_size, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, l, t + Inches(0.78), w, Inches(0.45),
             size=12, bold=False, color=lbl_color, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# Build charts (transparent bg for slides)
# ══════════════════════════════════════════════════════════════════════════

plt.rcParams.update({"font.family": "DejaVu Sans"})

# Chart A — Classifier comparison (balanced accuracy bar)
models_short = ["Logistic\nReg.", "Random\nForest", "SVM\n(RBF)", "Grad.\nBoosting", "PLS-DA"]
ba_vals = [baseline["models"][m]["mean_balanced_accuracy"]
           for m in ["LogisticRegression","RandomForest","SVM","GradientBoosting","PLSDA"]]
ba_std  = [baseline["models"][m]["std_balanced_accuracy"]
           for m in ["LogisticRegression","RandomForest","SVM","GradientBoosting","PLSDA"]]
oa_vals = [baseline["models"][m]["mean_overall_accuracy"]
           for m in ["LogisticRegression","RandomForest","SVM","GradientBoosting","PLSDA"]]

fig, ax = plt.subplots(figsize=(7, 3.5))
x = np.arange(5); w = 0.36
b1 = ax.bar(x-w/2, ba_vals, w, yerr=ba_std, capsize=4,
            color="#2E75B6", alpha=0.92, label="Balanced Accuracy", error_kw={"elinewidth":1.2})
b2 = ax.bar(x+w/2, oa_vals, w, color="#ED7D31", alpha=0.88, label="Overall Accuracy")
ax.axhline(1/3, color="#888", linestyle="--", lw=1, label="Chance (0.33)")
ax.set_xticks(x); ax.set_xticklabels(models_short, fontsize=10)
ax.set_ylim(0, 1); ax.set_ylabel("Score", fontsize=11)
ax.legend(fontsize=9, framealpha=0.5); ax.grid(axis="y", alpha=0.25)
ax.set_facecolor("none"); fig.patch.set_alpha(0)
for bar in b1:
    h_ = bar.get_height()
    ax.text(bar.get_x()+bar.get_width()/2, h_+0.018, f"{h_:.3f}",
            ha="center", va="bottom", fontsize=8.5, fontweight="bold", color="#1F497D")
plt.tight_layout(pad=0.4)
chartA = fig_to_stream(fig)

# Chart B — Per-site RF balanced accuracy
sites     = ["Bílý Kríž", "Hyytiälä", "Järvselja", "Lanzhot"]
rf_folds  = [1.000, 0.402, 0.583, 0.500]
colors_s  = ["#E41A1C","#377EB8","#4DAF4A","#984EA3"]
fig, ax = plt.subplots(figsize=(5.5, 3.4))
bars = ax.bar(sites, rf_folds, color=colors_s, alpha=0.88, edgecolor="white", lw=0.8)
ax.axhline(0.621, color="#333", linestyle="--", lw=1.4, label="Mean BA=0.621")
ax.axhline(1/3,   color="#888", linestyle=":", lw=1,   label="Chance")
ax.set_ylim(0, 1.18); ax.set_ylabel("Balanced Accuracy", fontsize=11)
ax.legend(fontsize=9, framealpha=0.4); ax.grid(axis="y", alpha=0.22)
ax.set_facecolor("none"); fig.patch.set_alpha(0)
for bar, v in zip(bars, rf_folds):
    ax.text(bar.get_x()+bar.get_width()/2, v+0.04, f"{v:.3f}",
            ha="center", fontsize=10.5, fontweight="bold", color="#1F497D")
plt.tight_layout(pad=0.4)
chartB = fig_to_stream(fig)

# Chart C — Confusion matrix
cm_raw  = np.array(baseline["models"]["RandomForest"]["total_confusion_matrix"])
row_sum = cm_raw.sum(axis=1, keepdims=True); row_sum[row_sum==0]=1
cm_norm = cm_raw/row_sum
cn = ["Conif.", "Broad.", "Mixed"]
fig, ax = plt.subplots(figsize=(4, 3.5))
im = ax.imshow(cm_norm, cmap="Blues", vmin=0, vmax=1)
plt.colorbar(im, ax=ax, fraction=0.048)
ax.set_xticks(range(3)); ax.set_yticks(range(3))
ax.set_xticklabels(cn, fontsize=10); ax.set_yticklabels(cn, fontsize=10)
ax.set_xlabel("Predicted", fontsize=10); ax.set_ylabel("True", fontsize=10)
for i in range(3):
    for j in range(3):
        v = cm_norm[i,j]; raw = cm_raw[i,j]
        c = "white" if v > 0.5 else "#222"
        ax.text(j, i, f"{v:.2f}\n(n={raw})", ha="center", va="center", fontsize=8.5, color=c)
fig.patch.set_alpha(0); ax.set_facecolor("none")
plt.tight_layout(pad=0.4)
chartC = fig_to_stream(fig)

# Chart D — Ablation bar chart
abl_labels = ["Baseline\nCASI 40b", "PCA-10", "ALS Only", "CASI+ALS\nFusion", "Binary\n(no mixed)", "Field\nSpectra", "Cross-\nSensor BK"]
abl_3cls   = [0.621, 0.600, 0.192, 0.663, None, 0.413, None]
abl_bin    = [0.932, None,  0.247, 0.932, 0.932, None, 0.875]
fig, ax = plt.subplots(figsize=(7.5, 3.5))
x = np.arange(7); w = 0.38
v3 = [v if v is not None else 0 for v in abl_3cls]
vb = [v if v is not None else 0 for v in abl_bin]
b3 = ax.bar(x-w/2, v3, w, color="#2E75B6", alpha=0.88, label="3-class BA")
bb = ax.bar(x+w/2, vb, w, color="#ED7D31", alpha=0.88, label="Binary BA")
for bar, a in zip(b3, [1 if v is not None else 0 for v in abl_3cls]):
    bar.set_alpha(a * 0.88 if a else 0)
for bar, a in zip(bb, [1 if v is not None else 0 for v in abl_bin]):
    bar.set_alpha(a * 0.88 if a else 0)
ax.axhline(1/3, color="#888", linestyle=":", lw=1)
ax.set_xticks(x); ax.set_xticklabels(abl_labels, fontsize=9)
ax.set_ylim(0, 1.08); ax.set_ylabel("Balanced Accuracy", fontsize=11)
ax.legend(fontsize=10, framealpha=0.4); ax.grid(axis="y", alpha=0.22)
ax.set_facecolor("none"); fig.patch.set_alpha(0)
for bar, v, vis in zip(b3, v3, abl_3cls):
    if vis is not None:
        ax.text(bar.get_x()+bar.get_width()/2, v+0.022, f"{v:.3f}", ha="center", fontsize=8, color="#1F497D")
for bar, v, vis in zip(bb, vb, abl_bin):
    if vis is not None:
        ax.text(bar.get_x()+bar.get_width()/2, v+0.022, f"{v:.3f}", ha="center", fontsize=8, color="#7B3F00")
plt.tight_layout(pad=0.4)
chartD = fig_to_stream(fig)

# Chart E — BK cross-sensor probabilities
bk_s    = [p["stand_id"].replace("BK_","") for p in cross["bk_predictions"]]
p_c     = [p["probabilities"]["coniferous"] for p in cross["bk_predictions"]]
p_b     = [p["probabilities"]["broadleaved"] for p in cross["bk_predictions"]]
p_m     = [p["probabilities"]["mixed"]       for p in cross["bk_predictions"]]
fig, ax = plt.subplots(figsize=(6, 3.4))
x5 = np.arange(len(bk_s)); w5 = 0.27
ax.bar(x5-w5, p_c, w5, label="P(coniferous)", color="#2E7D32", alpha=0.88)
ax.bar(x5,    p_m, w5, label="P(mixed)",      color="#7B1FA2", alpha=0.88)
ax.bar(x5+w5, p_b, w5, label="P(broadleaved)",color="#F57F17", alpha=0.88)
ax.set_xticks(x5); ax.set_xticklabels(bk_s, rotation=35, ha="right", fontsize=9)
ax.set_ylim(0, 1.08); ax.set_ylabel("Probability", fontsize=11)
ax.axhline(0.5, color="#888", linestyle="--", lw=0.8)
ax.legend(fontsize=8.5, framealpha=0.4); ax.grid(axis="y", alpha=0.22)
ax.set_facecolor("none"); fig.patch.set_alpha(0)
ax.annotate("Misclassified\n(27% broadleaved)", xy=(2, p_m[2]+0.02),
            xytext=(3.8, 0.7), arrowprops=dict(arrowstyle="->", color="red"),
            color="red", fontsize=8)
plt.tight_layout(pad=0.4)
chartE = fig_to_stream(fig)

# Chart F — Pipeline diagram
fig, ax = plt.subplots(figsize=(11, 2.5))
ax.set_xlim(0,11); ax.set_ylim(0,2.5); ax.axis("off")
fig.patch.set_alpha(0); ax.set_facecolor("none")
steps = [
    ("GeoTIFF\nTiles\n(58 stands)", 0.55, "#4E79A7"),
    ("Band\nSelection\n(40 bands)",  2.20, "#F28E2B"),
    ("Stand-level\nMean\nSpectrum",  3.85, "#E15759"),
    ("Normalise\n(per fold)",         5.50, "#76B7B2"),
    ("Classifier\n(RF / SVM…)",       7.15, "#59A14F"),
    ("LOSO-CV\n(4 folds)",            8.80, "#EDC948"),
    ("Metrics\n& Figures",           10.45, "#B07AA1"),
]
for label, x, col in steps:
    ax.add_patch(mpatches.FancyBboxPatch((x-0.68, 0.3), 1.36, 1.9,
                 boxstyle="round,pad=0.1", facecolor=col, edgecolor="white", lw=2, alpha=0.93))
    ax.text(x, 1.25, label, ha="center", va="center", fontsize=8.5,
            color="white", fontweight="bold", multialignment="center")
for i in range(len(steps)-1):
    x1 = steps[i][1]+0.68; x2 = steps[i+1][1]-0.68
    ax.annotate("", xy=(x2,1.25), xytext=(x1,1.25),
                arrowprops=dict(arrowstyle="->", color="#444", lw=2))
plt.tight_layout(pad=0.1)
chartF = fig_to_stream(fig)

# Chart G — 3-class vs binary delta
fig, ax = plt.subplots(figsize=(4.5, 3))
fold_sites = ["Bílý Kríž","Hyytiälä","Järvselja","Lanzhot","Mean"]
v3_ = [1.000, 0.402, 0.583, 0.500, 0.621]
vb_ = [1.000, 0.852, 0.875, 1.000, 0.932]
x_ = np.arange(5); w_ = 0.36
ax.bar(x_-w_/2, v3_, w_, color="#2E75B6", alpha=0.88, label="3-class")
ax.bar(x_+w_/2, vb_, w_, color="#ED7D31", alpha=0.88, label="Binary")
ax.set_xticks(x_); ax.set_xticklabels(fold_sites, rotation=20, ha="right", fontsize=9)
ax.set_ylim(0,1.15); ax.set_ylabel("Balanced Accuracy", fontsize=10)
ax.legend(fontsize=9, framealpha=0.4); ax.grid(axis="y", alpha=0.22)
ax.set_facecolor("none"); fig.patch.set_alpha(0)
plt.tight_layout(pad=0.4)
chartG = fig_to_stream(fig)

print("All charts generated.")


# ══════════════════════════════════════════════════════════════════════════
# Build Presentation
# ══════════════════════════════════════════════════════════════════════════
prs = new_prs()

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title slide
# ─────────────────────────────────────────────────────────────────────────
sld = blank_slide(prs)
fill_slide_bg(sld, DARK_BLUE)

# decorative accent strip bottom
add_rect(sld, 0, H - Inches(1.4), W, Inches(1.4), MED_BLUE)
add_rect(sld, 0, H - Inches(0.38), W, Inches(0.38), ACCENT)

add_text(sld, "Spectral Separability of Forest Types",
         Inches(0.7), Inches(0.9), Inches(12), Inches(1.2),
         size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sld, "Across European Biomes",
         Inches(0.7), Inches(2.05), Inches(12), Inches(0.9),
         size=36, bold=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_text(sld,
         "Airborne Hyperspectral Classification with Leave-One-Site-Out Cross-Validation",
         Inches(1.2), Inches(3.05), Inches(11), Inches(0.65),
         size=17, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, italic=True)

add_rect(sld, Inches(4.5), Inches(3.85), Inches(4.33), Inches(0.04), ACCENT)

add_text(sld, "Manan Aggarwal   |   Multi-disciplinary Hyperspectral Project   |   June 2026",
         Inches(1), Inches(4.1), Inches(11.33), Inches(0.55),
         size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# Stats strip
stat_box(sld, "58",      "Forest stands",    Inches(1.3),  Inches(5.25), bg=ACCENT)
stat_box(sld, "4",       "European sites",   Inches(3.9),  Inches(5.25), bg=rgb_hex(0x1A,0x62,0xA3))
stat_box(sld, "40",      "Spectral bands",   Inches(6.5),  Inches(5.25), bg=ACCENT)
stat_box(sld, "5",       "Classifiers",      Inches(9.1),  Inches(5.25), bg=rgb_hex(0x1A,0x62,0xA3))
stat_box(sld, "LOSO-CV","Spatial blocking", Inches(11.5), Inches(5.25), bg=ACCENT, val_size=18)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Problem + Dataset
# ─────────────────────────────────────────────────────────────────────────
sld = blank_slide(prs)
fill_slide_bg(sld, WHITE)
header_bar(sld, "The Problem & Dataset", "FREEDLES multi-site airborne hyperspectral survey")
slide_number(sld, 2)

# Left: research questions
bullet_box(sld,
    ["Can CASI-1500 (382–1052 nm) tell coniferous from broadleaved forest?",
     "Does it work across biomes — not just within a single site?",
     "Does ALS structural data improve classification?",
     "Do findings transfer to field spectroradiometer (350–2500 nm)?"],
    l=Inches(0.3), t=Inches(1.25), w=Inches(5.6), h=Inches(2.6),
    bg_color=LIGHT_GREY, size=13.5, title="Research Questions")

# Right: site table
col_header_bg = DARK_BLUE
table_data = [
    ("Hyytiälä (HY)",  "Finland",  "Boreal",           "28", "Pine, Spruce, Birch"),
    ("Järvselja (JS)", "Estonia",  "Hemiboreal",        "13", "Pine, Spruce, Birch, Alder"),
    ("Bílý Kríž (BK)","Czech R.", "Temperate montane", "7",  "Norway Spruce (plantation)"),
    ("Lanzhot (LZ)",   "Czech R.", "Temperate floodpl.","10", "Oak, Ash, Hornbeam"),
]
# draw table manually
col_x = [Inches(6.1), Inches(7.9), Inches(9.05), Inches(10.5), Inches(11.2)]
col_w = [Inches(1.75), Inches(1.1), Inches(1.4), Inches(0.65), Inches(1.9)]
row_h = Inches(0.52)
row_y0 = Inches(1.35)

headers = ["Site", "Country", "Biome", "n", "Species"]
for ci, (hdr, cw, cx) in enumerate(zip(headers, col_w, col_x)):
    add_rect(sld, cx, row_y0, cw, row_h, DARK_BLUE)
    add_text(sld, hdr, cx+Inches(0.05), row_y0+Inches(0.08), cw-Inches(0.05), row_h,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

for ri, row in enumerate(table_data):
    bg = LIGHT_GREY if ri % 2 == 0 else WHITE
    for ci, (val, cw, cx) in enumerate(zip(row, col_w, col_x)):
        add_rect(sld, cx, row_y0 + row_h*(ri+1), cw, row_h, bg,
                 line_color=rgb_hex(0xCC,0xCC,0xCC), line_width=0.5)
        add_text(sld, val,
                 cx+Inches(0.05), row_y0+row_h*(ri+1)+Inches(0.07),
                 cw-Inches(0.07), row_h,
                 size=10, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT)

# Class distribution heatmap (inline)
counts = [[22,4,2],[4,4,5],[7,0,0],[0,7,3]]
site_short = ["HY","JS","BK","LZ"]
cls_short  = ["Conif.","Broad.","Mixed"]
fig_h, ax_h = plt.subplots(figsize=(3.2, 2.8))
im_h = ax_h.imshow(counts, cmap="YlOrRd", aspect="auto", vmin=0, vmax=22)
ax_h.set_xticks(range(3)); ax_h.set_yticks(range(4))
ax_h.set_xticklabels(cls_short, fontsize=10); ax_h.set_yticklabels(site_short, fontsize=10)
ax_h.set_title("Stand counts by site & type", fontsize=10)
for i in range(4):
    for j in range(3):
        v = counts[i][j]
        ax_h.text(j, i, str(v), ha="center", va="center", fontsize=13,
                  fontweight="bold", color="white" if v>10 else "black")
fig_h.patch.set_alpha(0); ax_h.set_facecolor("none")
plt.tight_layout(pad=0.3)
hmap_stream = fig_to_stream(fig_h)
add_image_stream(sld, hmap_stream, Inches(6.1), Inches(4.05), Inches(3.8))

# Key fact box
add_rect(sld, Inches(0.3), Inches(3.95), Inches(5.6), Inches(1.25), rgb_hex(0x1A,0x62,0xA3))
add_text(sld, "Evaluation: Leave-One-Site-Out (LOSO) cross-validation\n"
              "Train on 3 sites → test on the 4th (strict spatial blocking)\n"
              "Primary metric: Balanced Accuracy (handles class imbalance)",
         Inches(0.5), Inches(4.0), Inches(5.3), Inches(1.15),
         size=12.5, color=WHITE)

# Caption
add_text(sld, "57% coniferous  |  26% broadleaved  |  17% mixed",
         Inches(0.3), Inches(5.28), Inches(5.6), Inches(0.4),
         size=12, color=DARK_GREY, italic=True)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Pipeline
# ─────────────────────────────────────────────────────────────────────────
sld = blank_slide(prs)
fill_slide_bg(sld, WHITE)
header_bar(sld, "Classification Pipeline", "From raw GeoTIFF tiles to evaluated metrics")
slide_number(sld, 3)

add_image_stream(sld, chartF, Inches(0.5), Inches(1.4), Inches(12.2))

# Key decisions boxes
boxes = [
    ("Feature Extraction",
     ["Stand-level MEAN reflectance per tile",
      "All valid pixels averaged band-wise",
      "Avoids overfitting (n=58 small sample)"],
     Inches(0.3), Inches(4.0)),
    ("Band Selection",
     ["48 CASI bands → 40 valid",
      "Remove 8 water-vapour bands",
      "(895-1003, 1092-1168, 1302-1528, 1737-2038 nm)"],
     Inches(4.55), Inches(4.0)),
    ("Normalisation",
     ["z-score fit on TRAINING fold only",
      "Applied to test fold → no data leakage",
      "Refit fresh in each of 4 LOSO folds"],
     Inches(8.8), Inches(4.0)),
]
for title, items, lx, ty in boxes:
    bullet_box(sld, items, lx, ty, Inches(4.05), Inches(2.15),
               bg_color=LIGHT_GREY, size=12.5, title=title)

add_text(sld, "Critical data fixes discovered during development",
         Inches(0.3), Inches(6.25), Inches(8), Inches(0.4),
         size=13, bold=True, color=DARK_BLUE)
add_text(sld, "nodata=10000 (not 0)  •  scale factor /10000  •  wavelengths from ENVI headers  •  ALS veg code=4 at BK/LZ",
         Inches(0.3), Inches(6.6), Inches(12.5), Inches(0.5),
         size=11.5, color=DARK_GREY, italic=True)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Classifier Results
# ─────────────────────────────────────────────────────────────────────────
sld = blank_slide(prs)
fill_slide_bg(sld, WHITE)
header_bar(sld, "Classifier Comparison — 3-class LOSO-CV",
           "58 stands, 4 folds, 3 forest types (coniferous / broadleaved / mixed)")
slide_number(sld, 4)

add_image_stream(sld, chartA, Inches(0.25), Inches(1.3), Inches(7.2))

# Summary table on right
tbl_data = [
    ("Logistic Reg.",  "0.542", "0.698", "0.470"),
    ("Random Forest",  "0.621", "0.756", "0.576"),
    ("SVM (RBF)",      "0.592", "0.694", "0.438"),
    ("Grad. Boosting", "0.592", "0.706", "0.555"),
    ("PLS-DA",         "0.591", "—",     "—"    ),
]
hdr_y = Inches(1.35); rx = Inches(7.65); rw = Inches(1.45)
for ci, (hdr, cw) in enumerate(zip(["Model","BA","OA","F1"],
                                     [Inches(2.3),Inches(0.95),Inches(0.95),Inches(0.95)])):
    add_rect(sld, rx + sum([Inches(2.3),Inches(0.95),Inches(0.95),Inches(0.95)][:ci]),
             hdr_y, cw, Inches(0.5), DARK_BLUE)
    add_text(sld, hdr,
             rx + sum([Inches(2.3),Inches(0.95),Inches(0.95),Inches(0.95)][:ci]) + Inches(0.05),
             hdr_y+Inches(0.07), cw, Inches(0.5),
             size=11, bold=True, color=WHITE)
col_xs = [rx, rx+Inches(2.3), rx+Inches(3.25), rx+Inches(4.2)]
col_ws = [Inches(2.3), Inches(0.95), Inches(0.95), Inches(0.95)]
for ri, row in enumerate(tbl_data):
    is_best = ri == 1
    bg = rgb_hex(0xBD,0xD7,0xEE) if is_best else (LIGHT_GREY if ri%2==0 else WHITE)
    for ci, (val, cx, cw) in enumerate(zip(row, col_xs, col_ws)):
        add_rect(sld, cx, hdr_y+Inches(0.5)*(ri+1), cw, Inches(0.5), bg,
                 line_color=rgb_hex(0xCC,0xCC,0xCC), line_width=0.4)
        add_text(sld, val,
                 cx+Inches(0.06), hdr_y+Inches(0.5)*(ri+1)+Inches(0.08),
                 cw, Inches(0.5),
                 size=11, bold=(is_best and ci > 0), color=DARK_GREY)

add_text(sld, "★ Random Forest wins on BA and Macro-F1",
         Inches(7.65), hdr_y+Inches(3.2), Inches(5.3), Inches(0.5),
         size=12.5, bold=True, color=DARK_BLUE)

bullet_box(sld,
    ["Broadleaved recall ≈ 0.94 → clear spectral signature",
     "Mixed recall ≈ 0.00 → completely undetected",
     "Site shortcut ratio 1.14× → model uses biology, not site offsets"],
    Inches(7.65), hdr_y+Inches(3.7), Inches(5.3), Inches(2.3),
    bg_color=rgb_hex(0xF0,0xF7,0xFF), size=12.5)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Per-site + Confusion matrix
# ─────────────────────────────────────────────────────────────────────────
sld = blank_slide(prs)
fill_slide_bg(sld, WHITE)
header_bar(sld, "Per-Site Performance & Error Analysis",
           "Where does the model succeed — and where does it fail?")
slide_number(sld, 5)

add_image_stream(sld, chartB, Inches(0.25), Inches(1.25), Inches(6.2))
add_image_stream(sld, chartC, Inches(6.6),  Inches(1.15), Inches(4.5))

# Insight boxes below
add_rect(sld, Inches(0.3), Inches(5.2), Inches(5.9), Inches(1.05), rgb_hex(0xE8,0xF4,0xE8))
add_text(sld, "Bílý Kríž BA=1.0 → trivial (100% coniferous test set). "
              "Remove it: mean BA over HY+JS+LZ = 0.495",
         Inches(0.45), Inches(5.25), Inches(5.65), Inches(0.95),
         size=12.5, color=GREEN)

add_rect(sld, Inches(6.5), Inches(5.2), Inches(6.5), Inches(1.05), rgb_hex(0xFF,0xED,0xED))
add_text(sld, "Mixed class = 0/10 correct. All misclassified as CONIFEROUS. "
              "Stand-level mean spectra can't resolve spatial mixtures → needs crown-level analysis.",
         Inches(6.65), Inches(5.25), Inches(6.2), Inches(0.95),
         size=12.5, color=rgb_hex(0xC0,0x00,0x00))

add_text(sld, "Hyytiälä (BA=0.402) is the hardest fold — all 3 classes present, 28 stands",
         Inches(0.3), Inches(6.35), Inches(12.5), Inches(0.45),
         size=12, color=DARK_GREY, italic=True)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Binary breakthrough
# ─────────────────────────────────────────────────────────────────────────
sld = blank_slide(prs)
fill_slide_bg(sld, WHITE)
header_bar(sld, "Key Finding: Binary Classification",
           "Drop the mixed class — conifer vs. broadleaf is biome-stable")
slide_number(sld, 6)

add_image_stream(sld, chartG, Inches(0.3), Inches(1.3), Inches(5.6))

# Big stats
stat_box(sld, "0.621", "3-class BA",   Inches(6.4), Inches(1.5),
         bg=rgb_hex(0x9B,0xB9,0xD6), val_color=DARK_BLUE, lbl_color=DARK_BLUE, val_size=28)
add_text(sld, "→", Inches(8.95), Inches(1.7), Inches(0.7), Inches(0.9),
         size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
stat_box(sld, "0.932", "Binary BA",    Inches(9.75), Inches(1.5),
         bg=GREEN, val_color=WHITE, lbl_color=LIGHT_BLUE, val_size=28)

add_text(sld, "+31.1 percentage points",
         Inches(6.4), Inches(2.95), Inches(6.7), Inches(0.55),
         size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

bullet_box(sld,
    ["Hyytiälä: 0.402 → 0.852  (+45 pp)",
     "Järvselja: 0.583 → 0.875  (+29 pp)",
     "Lanzhot:   0.500 → 1.000  (+50 pp)",
     "Bílý Kríž: 1.000 → 1.000  (±0 pp)"],
    Inches(6.4), Inches(3.6), Inches(6.6), Inches(2.2),
    bg_color=LIGHT_GREY, size=13, title="Per-fold gain from going binary")

add_rect(sld, Inches(0.3), Inches(5.4), Inches(5.6), Inches(1.7),
         rgb_hex(0x1A,0x62,0xA3))
add_text(sld,
         "Interpretation:\nThe biome-generalisation problem is SOLVED for conifer vs. broadleaf.\n"
         "The residual 3-class error is entirely the mixed label — not sensor limitations.",
         Inches(0.5), Inches(5.45), Inches(5.3), Inches(1.6),
         size=13, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Ablation study
# ─────────────────────────────────────────────────────────────────────────
sld = blank_slide(prs)
fill_slide_bg(sld, WHITE)
header_bar(sld, "Ablation Study — What Helps and What Doesn't",
           "Systematic isolation of each design decision")
slide_number(sld, 7)

add_image_stream(sld, chartD, Inches(0.25), Inches(1.25), Inches(8.0))

abl_table = [
    ("PCA-10",        "0.600", "−0.021", "PCA discards biologically relevant low-variance bands"),
    ("ALS only",      "0.192", "−0.429", "Height/cover not diagnostic across biomes"),
    ("CASI+ALS",      "0.663", "+0.042", "Modest 3-class gain; zero binary gain"),
    ("Binary",        "0.932", "+0.311", "Largest gain — mixed class was the bottleneck"),
    ("Field spectra", "0.413", "−0.208", "BK fold collapses (training imbalance)"),
    ("Cross-sensor",  "0.875", "+—",     "7/8 BK stands correct via spectral resampling"),
]
rx2 = Inches(8.4); ry0 = Inches(1.4)
hdrs = ["Ablation","3-cls BA","Δ","Interpretation"]
hcws = [Inches(1.4),Inches(1.0),Inches(0.85),Inches(2.7)]
hcxs = [rx2, rx2+Inches(1.4), rx2+Inches(2.4), rx2+Inches(3.25)]
for ci, (h, cw, cx) in enumerate(zip(hdrs, hcws, hcxs)):
    add_rect(sld, cx, ry0, cw, Inches(0.46), DARK_BLUE)
    add_text(sld, h, cx+Inches(0.05), ry0+Inches(0.07), cw, Inches(0.46),
             size=10, bold=True, color=WHITE)
for ri, row in enumerate(abl_table):
    bg = LIGHT_GREY if ri%2==0 else WHITE
    is_best = ri == 3
    if is_best: bg = rgb_hex(0xE2,0xF0,0xD9)
    for ci, (val, cw, cx) in enumerate(zip(row, hcws, hcxs)):
        add_rect(sld, cx, ry0+Inches(0.46)*(ri+1), cw, Inches(0.46), bg,
                 line_color=rgb_hex(0xCC,0xCC,0xCC), line_width=0.4)
        delta_col = GREEN if (ci==2 and "+" in str(val)) else (rgb_hex(0xC0,0,0) if (ci==2 and "−" in str(val)) else DARK_GREY)
        add_text(sld, val, cx+Inches(0.05), ry0+Inches(0.46)*(ri+1)+Inches(0.07),
                 cw, Inches(0.46), size=9.5, bold=(is_best and ci<2),
                 color=delta_col if ci==2 else DARK_GREY)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Field Spectra & Cross-sensor
# ─────────────────────────────────────────────────────────────────────────
sld = blank_slide(prs)
fill_slide_bg(sld, WHITE)
header_bar(sld, "BilyKriz Field Spectra & Cross-Sensor Transfer",
           "350–2500 nm spectroradiometer + CASI model applied to field data")
slide_number(sld, 8)

add_image_stream(sld, chartE, Inches(0.25), Inches(1.3), Inches(6.0))

# Separability figure
sep_fig_path = FIG_DIR / "field_spectra_separability.png"
if sep_fig_path.exists():
    add_image_path(sld, sep_fig_path, Inches(6.4), Inches(1.3), Inches(6.6))

# Key findings
bullet_box(sld,
    ["Field LOSO-CV (31 stands, 3 sites): BA = 0.413",
     "BK fold collapses: 19 broadleaved vs 4 coniferous in training",
     "Top separable region: SWIR-2 (2317–2490 nm) — cellulose/lignin",
     "SWIR not in CASI range yet CASI binary BA = 0.932 → VIS-NIR sufficient",
     "Cross-sensor BK: 7/8 correct (BA = 0.875), misclassified Spruce3 has 27% broadleaved"],
    Inches(0.25), Inches(5.1), Inches(12.8), Inches(2.05),
    bg_color=rgb_hex(0xF0,0xF5,0xFF), size=12.5,
    title="Key findings — Field Spectroscopy")


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Iteration log (key failures & fixes)
# ─────────────────────────────────────────────────────────────────────────
sld = blank_slide(prs)
fill_slide_bg(sld, WHITE)
header_bar(sld, "Iteration Log — What We Fixed Along the Way",
           "8 non-trivial bugs discovered and resolved during pipeline development")
slide_number(sld, 9)

failures = [
    ("nodata = 10000, not 0",
     "Pixel mask let 10000 DN through; post-scale = 1.0 reflectance corruption"),
    ("Reflectance scale /10000",
     "Raw DN range 0–2350; forgot to divide by 10000 — band importance was raw DN"),
    ("Wavelengths from ENVI headers",
     "GeoTIFFs have no spectral metadata; water-vapour exclusion silently skipped"),
    ("Stand ID suffix stripping",
     "HY_PINE1_CASI.tif → stand_id=HY_PINE1_CASI; no metadata match"),
    ("ALS veg code = 4 (not 1)",
     "BK/LZ used cls=4 for medium veg; cls==1 filter gave zero points"),
    ("Single-class fold crash",
     "BK fold = only coniferous; sklearn classification_report raised ValueError"),
    ("SVM multi-class solver",
     "liblinear incompatible with balanced multi-class; switched to lbfgs"),
    ("f-string backslash Python 3.12",
     "f'{f[\"ba\"]}' is SyntaxError in 3.12; pre-computed variable needed"),
]

# Two columns
col1 = failures[:4]; col2 = failures[4:]
for ci, col in enumerate([col1, col2]):
    for ri, (bug, fix) in enumerate(col):
        lx = Inches(0.3) if ci==0 else Inches(6.7)
        ty = Inches(1.4 + ri * 1.35)
        add_rect(sld, lx, ty, Inches(6.2), Inches(1.2), LIGHT_GREY,
                 line_color=rgb_hex(0xBB,0xBB,0xBB), line_width=0.5)
        add_rect(sld, lx, ty, Inches(0.22), Inches(1.2), ACCENT)
        add_text(sld, bug, lx+Inches(0.3), ty+Inches(0.06), Inches(5.8), Inches(0.45),
                 size=12, bold=True, color=DARK_BLUE)
        add_text(sld, fix, lx+Inches(0.3), ty+Inches(0.5), Inches(5.8), Inches(0.62),
                 size=11, color=DARK_GREY)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Conclusions & Next Steps
# ─────────────────────────────────────────────────────────────────────────
sld = blank_slide(prs)
fill_slide_bg(sld, WHITE)
header_bar(sld, "Conclusions & Next Steps", "")
slide_number(sld, 10)

# Left — conclusions
add_rect(sld, Inches(0.3), Inches(1.3), Inches(6.2), Inches(5.6), DARK_BLUE)
add_text(sld, "Conclusions", Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.55),
         size=18, bold=True, color=WHITE)
concl = [
    "Binary conifer/broadleaf: BA = 0.932\nBiome generalisation is SOLVED",
    "3-class problem: BA = 0.621\nMixed class is entirely undetected",
    "ALS gives +4.1 pp for 3-class\nZero benefit for binary",
    "SWIR-2 most separable (field data)\nBut VIS-NIR already sufficient",
    "Cross-sensor transfer: 7/8 BK\nBA = 0.875 via spectral resampling",
]
icons = ["✓", "△", "⊕", "★", "↔"]
for i, (txt, icon) in enumerate(zip(concl, icons)):
    ty = Inches(1.97 + i*1.0)
    add_rect(sld, Inches(0.35), ty, Inches(0.65), Inches(0.85), ACCENT)
    add_text(sld, icon, Inches(0.37), ty+Inches(0.08), Inches(0.6), Inches(0.75),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sld, txt, Inches(1.08), ty+Inches(0.05), Inches(5.25), Inches(0.88),
             size=12.5, color=WHITE)

# Right — next steps
add_rect(sld, Inches(6.7), Inches(1.3), Inches(6.3), Inches(2.65), LIGHT_GREY)
add_text(sld, "What's Next", Inches(6.9), Inches(1.4), Inches(5.9), Inches(0.55),
         size=16, bold=True, color=DARK_BLUE)
next_steps = [
    "Crown-level / patch-level features for mixed class",
    "Species-level within coniferous type (Pine vs. Spruce)",
    "Multi-date imagery — phenological separability",
    "Probability calibration for operational deployment",
]
for i, ns in enumerate(next_steps):
    add_text(sld, f"{i+1}.  {ns}",
             Inches(6.9), Inches(1.9 + i*0.6), Inches(5.9), Inches(0.55),
             size=12.5, color=DARK_GREY)

# Right — key numbers summary
add_rect(sld, Inches(6.7), Inches(4.1), Inches(6.3), Inches(2.8), rgb_hex(0x1A,0x62,0xA3))
add_text(sld, "Key Numbers at a Glance",
         Inches(6.9), Inches(4.18), Inches(5.9), Inches(0.5),
         size=14, bold=True, color=WHITE)
kn = [
    ("0.621", "3-class Balanced Acc."),
    ("0.932", "Binary Balanced Acc."),
    ("0.663", "CASI + ALS fusion"),
    ("0.875", "Cross-sensor BK"),
    ("2317nm","Most separable band"),
    ("7/8",   "BK stands correct"),
]
for i, (val, lbl) in enumerate(kn):
    col = i % 3; row = i // 3
    lx = Inches(6.85) + col * Inches(2.1)
    ty = Inches(4.72) + row * Inches(1.0)
    add_rect(sld, lx, ty, Inches(2.0), Inches(0.85), rgb_hex(0x2E,0x75,0xB6))
    add_text(sld, val, lx, ty+Inches(0.04), Inches(2.0), Inches(0.42),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sld, lbl, lx, ty+Inches(0.46), Inches(2.0), Inches(0.38),
             size=9.5, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────
prs.save(str(PPTX_OUT))
print(f"\nPresentation saved: {PPTX_OUT}")
print(f"File size: {PPTX_OUT.stat().st_size / 1024:.1f} KB")
print(f"Slides:    {len(prs.slides)}")
